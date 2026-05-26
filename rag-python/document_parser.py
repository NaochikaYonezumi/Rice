import io
import os
import csv
import logging
from typing import Optional

logger = logging.getLogger(__name__)


class UnsupportedFileError(Exception):
    """OCR 等が必要でテキスト抽出できない場合に投げる。

    main.py 側で catch して 422 + "ocr_not_supported" コードを返す。
    """
    def __init__(self, message: str, kind: str = "ocr_not_supported"):
        super().__init__(message)
        self.kind = kind


# OCR 必須となる可能性の高い画像系拡張子
_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp", ".heic", ".heif")

# OCR 言語の組み合わせ。複数指定時は tesseract が両方の学習データを使う。
# 環境変数 OCR_LANG で上書き可 (例: "jpn+eng", "jpn_vert+jpn+eng")。
_OCR_LANG = os.environ.get("OCR_LANG", "jpn+eng")


def parse_file(content: bytes, filename: str, mime_type: str) -> Optional[str]:
    name = filename.lower()

    # 画像ファイル: Tesseract OCR で取り込む。
    # 失敗時は UnsupportedFileError を投げる (空文字列でも空 = 抽出失敗扱い)。
    if name.endswith(_IMAGE_EXTS) or mime_type.startswith("image/"):
        return _parse_image(content, filename)

    if name.endswith(".pdf") or "pdf" in mime_type:
        return _parse_pdf(content, filename)

    if name.endswith(".docx") or name.endswith(".doc") or "wordprocessingml" in mime_type or mime_type == "application/msword":
        return _parse_docx(content)

    if name.endswith(".pptx") or name.endswith(".ppt") or "presentationml" in mime_type:
        return _parse_pptx(content)

    if name.endswith(".xlsx") or "spreadsheetml" in mime_type:
        return _parse_xlsx(content)

    if name.endswith(".xls") or mime_type == "application/vnd.ms-excel":
        return _parse_xls(content)

    if name.endswith(".csv") or mime_type == "text/csv":
        return _parse_csv(content)

    if name.endswith(".odt") or "opendocument.text" in mime_type:
        return _parse_odt(content)

    if name.endswith(".ods") or "opendocument.spreadsheet" in mime_type:
        return _parse_ods(content)

    if name.endswith(".odp") or "opendocument.presentation" in mime_type:
        return _parse_odp(content)

    if name.endswith(".rtf") or mime_type == "application/rtf" or mime_type == "text/rtf":
        return _parse_rtf(content)

    if name.endswith(".epub") or mime_type == "application/epub+zip":
        return _parse_epub(content)

    if name.endswith((".html", ".htm")) or "html" in mime_type:
        return _parse_html(content)

    if name.endswith((".md", ".txt", ".rst", ".log", ".json", ".yaml", ".yml", ".xml", ".toml", ".ini", ".cfg")) or mime_type.startswith("text/"):
        return _decode_text(content)

    # 不明な形式はテキストとして試みる
    result = _decode_text(content)
    if result and len(result.strip()) > 50:
        return result

    logger.warning(f"Unsupported file type: {mime_type} / {filename}")
    return None


def _decode_text(content: bytes) -> Optional[str]:
    for enc in ("utf-8", "utf-8-sig", "shift_jis", "cp932", "euc_jp", "latin-1"):
        try:
            return content.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return None


def _normalize_cjk(text: str) -> str:
    """PDFフォントの誤エンコーディングで混入する韓国語文字を日本語に修正する"""
    import unicodedata
    # よく誤変換される韓国語→日本語カタカナの対応表
    _hangul_to_katakana = {
        '\uC18C': 'ソ', '\uC2A4': 'ス', '\uD2B8': 'ト', '\uC6CC': 'ウ',
        '\uD06C': 'ク', '\uC138': 'セ', '\uB85C': 'ロ', '\uB2C8': 'ニ',
        '\uCE90': 'ケ', '\uD14C': 'テ', '\uC5D0': 'エ', '\uC774': 'イ',
        '\uC544': 'ア', '\uCE74': 'カ', '\uC0AC': 'サ', '\uB098': 'ナ',
        '\uD558': 'ハ', '\uB9C8': 'マ', '\uC57C': 'ヤ', '\uB77C': 'ラ',
        '\uC640': 'ワ', '\uC624': 'オ', '\uCF54': 'コ', '\uC2DC': 'シ',
        '\uCE58': 'チ', '\uD754': 'フ', '\uD5E4': 'ヘ', '\uBE44': 'ビ',
        '\uD53C': 'ピ', '\uBBF8': 'ミ', '\uB9AC': 'リ',
    }
    result = []
    for ch in text:
        result.append(_hangul_to_katakana.get(ch, ch))
    return "".join(result)


def _parse_pdf(content: bytes, filename: str = "PDF") -> Optional[str]:
    # PyMuPDF を優先（日本語フォント処理が良好）
    try:
        import fitz
        doc = fitz.open(stream=content, filetype="pdf")
        pages = [page.get_text() for page in doc]
        text = "\n\n".join(p for p in pages if p.strip())
        if text.strip():
            return _normalize_cjk(text)
    except Exception as e:
        logger.warning(f"PyMuPDF failed: {e}")

    # フォールバック: pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        text = "\n\n".join(p for p in pages if p.strip())
        if text.strip():
            return _normalize_cjk(text)
    except Exception as e:
        logger.error(f"PDF parse failed: {e}")

    # ここまで来た = PDF からテキストを抽出できなかった (= スキャン画像 PDF など)。
    # OCR フォールバック: pdf2image で 1 ページずつ画像化して Tesseract に通す。
    ocr_text = _ocr_pdf_pages(content, filename)
    if ocr_text and ocr_text.strip():
        return _normalize_cjk(ocr_text)

    raise UnsupportedFileError(
        f"PDF ({filename}) からテキストを抽出できませんでした。"
        " 通常の抽出 + OCR どちらも失敗しています。画質が極端に低い、"
        " もしくは PDF が壊れている可能性があります。"
    )


def _parse_image(content: bytes, filename: str = "image") -> Optional[str]:
    """画像 1 枚を Tesseract OCR に通してテキストを返す。"""
    try:
        import pytesseract
        from PIL import Image, ImageOps
    except Exception as e:  # pragma: no cover
        raise UnsupportedFileError(
            f"OCR ライブラリが利用できません: {e}",
            kind="ocr_unavailable",
        )

    try:
        img = Image.open(io.BytesIO(content))
        # EXIF の orientation を反映してスキャン画像の向きを正す
        img = ImageOps.exif_transpose(img)
        # アルファ付き画像はグレースケール化前に背景白で合成 (Tesseract が透明部を黒として誤認するため)
        if img.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1])
            img = bg
        else:
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img, lang=_OCR_LANG)
        text = (text or "").strip()
        if text:
            return _normalize_cjk(text)
    except UnsupportedFileError:
        raise
    except Exception as e:
        logger.warning(f"Image OCR failed for {filename}: {e}")

    raise UnsupportedFileError(
        f"画像 ({filename}) から文字を読み取れませんでした。"
        " 文字がぼけている / 解像度が低い / 文字が無い 画像の可能性があります。",
        kind="ocr_empty",
    )


def _ocr_pdf_pages(content: bytes, filename: str = "PDF") -> Optional[str]:
    """PDF をページ単位で画像化 → Tesseract OCR で連結テキストを返す。"""
    try:
        import pytesseract
        from pdf2image import convert_from_bytes
    except Exception as e:
        logger.warning(f"OCR libs not available for PDF fallback: {e}")
        return None

    try:
        # 200 DPI で十分。多ページ PDF を OCR にかけるのは重いので、
        # 環境変数 OCR_PDF_MAX_PAGES (デフォルト 30) でハードリミットを設ける。
        max_pages = int(os.environ.get("OCR_PDF_MAX_PAGES", "30"))
        images = convert_from_bytes(content, dpi=200, fmt="png")
        if not images:
            return None
        parts: list = []
        for i, img in enumerate(images[:max_pages], 1):
            try:
                page_text = pytesseract.image_to_string(img, lang=_OCR_LANG)
                page_text = (page_text or "").strip()
                if page_text:
                    parts.append(f"[ページ {i}]\n{page_text}")
            except Exception as e:
                logger.warning(f"OCR page {i} failed for {filename}: {e}")
        if len(images) > max_pages:
            parts.append(
                f"[注意] ページ数が {max_pages} を超えるため、{max_pages + 1} ページ以降は OCR を省略しました。"
            )
        return "\n\n".join(parts) if parts else None
    except Exception as e:
        logger.error(f"PDF OCR fallback failed for {filename}: {e}")
        return None


def _parse_docx(content: bytes) -> Optional[str]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"DOCX parse failed: {e}")
        return None


def _extract_text_from_pptx_shape(shape, out: list) -> None:
    """1 つのシェイプから抽出可能なテキストをすべて out に追記する。

    対象: テキストフレーム / 表 / グループシェイプ (再帰) / 代替テキスト。
    pptx は shape の種類が多く、`has_text_frame` だけだとテーブル中心や
    SmartArt 中心のスライドが空判定 → 422 parse_empty になる。
    """
    try:
        # テキストフレーム (タイトル / 段落)
        if getattr(shape, "has_text_frame", False):
            for para in shape.text_frame.paragraphs:
                t = para.text.strip() if para.text else ""
                if t:
                    out.append(t)
        # 表
        if getattr(shape, "has_table", False):
            try:
                for row in shape.table.rows:
                    cells = []
                    for cell in row.cells:
                        ct = (cell.text or "").strip()
                        if ct:
                            cells.append(ct)
                    if cells:
                        out.append("\t".join(cells))
            except Exception:
                pass
        # グループシェイプ (再帰)
        if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP = 6
            try:
                for sub in shape.shapes:
                    _extract_text_from_pptx_shape(sub, out)
            except Exception:
                pass
        # 代替テキスト (画像説明など)
        try:
            alt = (getattr(shape, "alternative_text", "") or "").strip()
            if alt and len(alt) > 1:
                out.append(f"[alt] {alt}")
        except Exception:
            pass
    except Exception:
        # 個別 shape のパース失敗は無視して全体は続行
        pass


def _parse_pptx(content: bytes) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list = []
            for shape in slide.shapes:
                _extract_text_from_pptx_shape(shape, slide_texts)
            # 発表者ノート
            try:
                if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                    notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if notes_text:
                        slide_texts.append(f"[ノート]\n{notes_text}")
            except Exception:
                pass
            if slide_texts:
                # 重複行を保持しつつ純粋な空行を除去
                non_empty = [s for s in slide_texts if s and s.strip()]
                if non_empty:
                    parts.append(f"[スライド {i}]\n" + "\n".join(non_empty))
        return "\n\n".join(parts) if parts else None
    except Exception as e:
        logger.error(f"PPTX parse failed: {e}")
        return None


def _parse_xlsx(content: bytes) -> Optional[str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"[シート: {sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"XLSX parse failed: {e}")
        return None


def _parse_xls(content: bytes) -> Optional[str]:
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=content)
        parts = []
        for sheet in wb.sheets():
            rows = []
            for i in range(sheet.nrows):
                cells = [str(sheet.cell_value(i, j)) for j in range(sheet.ncols)]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"[シート: {sheet.name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"XLS parse failed: {e}")
        return None


def _parse_csv(content: bytes) -> Optional[str]:
    text = _decode_text(content)
    if not text:
        return None
    try:
        reader = csv.reader(io.StringIO(text))
        rows = ["\t".join(row) for row in reader if any(c.strip() for c in row)]
        return "\n".join(rows)
    except Exception as e:
        logger.error(f"CSV parse failed: {e}")
        return text


def _parse_odt(content: bytes) -> Optional[str]:
    try:
        from odf.opendocument import load
        from odf.text import P
        doc = load(io.BytesIO(content))
        parts = [str(p) for p in doc.getElementsByType(P) if str(p).strip()]
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"ODT parse failed: {e}")
        return None


def _parse_ods(content: bytes) -> Optional[str]:
    try:
        from odf.opendocument import load
        from odf.table import Table, TableRow, TableCell
        from odf.text import P
        doc = load(io.BytesIO(content))
        parts = []
        for table in doc.getElementsByType(Table):
            rows = []
            for row in table.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    ps = cell.getElementsByType(P)
                    cells.append(" ".join(str(p) for p in ps))
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append("\n".join(rows))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"ODS parse failed: {e}")
        return None


def _parse_odp(content: bytes) -> Optional[str]:
    try:
        from odf.opendocument import load
        from odf.text import P
        doc = load(io.BytesIO(content))
        parts = [str(p) for p in doc.getElementsByType(P) if str(p).strip()]
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"ODP parse failed: {e}")
        return None


def _parse_rtf(content: bytes) -> Optional[str]:
    try:
        from striprtf.striprtf import rtf_to_text
        text = _decode_text(content)
        if text:
            return rtf_to_text(text)
    except Exception as e:
        logger.error(f"RTF parse failed: {e}")
    return None


def _parse_epub(content: bytes) -> Optional[str]:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(io.BytesIO(content))
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n").strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"EPUB parse failed: {e}")
        return None


def _parse_html(content: bytes) -> Optional[str]:
    try:
        from bs4 import BeautifulSoup
        text = _decode_text(content) or content.decode("utf-8", errors="replace")
        soup = BeautifulSoup(text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n").strip()
    except Exception as e:
        logger.error(f"HTML parse failed: {e}")
        return None
